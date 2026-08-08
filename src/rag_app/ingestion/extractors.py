from __future__ import annotations

import csv
import io
from collections.abc import Callable
from dataclasses import dataclass
from pathlib import Path
from urllib.parse import urlsplit


class UnsupportedDocumentError(ValueError):
    """Формат документа не поддерживается безопасным extractor-ом."""


class EmptyDocumentError(ValueError):
    """Из документа не удалось извлечь непустой текст."""


@dataclass(frozen=True, slots=True)
class ExtractedUnit:
    text: str
    location: str


@dataclass(frozen=True, slots=True)
class DocumentLink:
    label: str
    target: str


@dataclass(frozen=True, slots=True)
class ExtractedDocument:
    title: str
    units: tuple[ExtractedUnit, ...]
    links: tuple[DocumentLink, ...] = ()


def _decode(payload: bytes) -> str:
    for encoding in ("utf-8-sig", "cp1251"):
        try:
            return payload.decode(encoding)
        except UnicodeDecodeError:
            continue
    raise ValueError("Текстовый файл должен быть в UTF-8 или Windows-1251")


def _text(path: Path) -> list[ExtractedUnit]:
    return [ExtractedUnit(_decode(path.read_bytes()), "текст")]


def _csv(path: Path) -> list[ExtractedUnit]:
    content = _decode(path.read_bytes())
    sample = content[:4096]
    try:
        dialect = csv.Sniffer().sniff(sample, delimiters=",;\t")
    except csv.Error:
        dialect = csv.excel
    rows = [
        " | ".join(cell.strip() for cell in row)
        for row in csv.reader(io.StringIO(content), dialect)
    ]
    return [ExtractedUnit("\n".join(row for row in rows if row.strip(" |")), "таблица")]


def _pdf(path: Path) -> list[ExtractedUnit]:
    from pypdf import PdfReader

    reader = PdfReader(path)
    return [
        ExtractedUnit(text, f"стр. {number}")
        for number, page in enumerate(reader.pages, start=1)
        if (text := (page.extract_text() or "").strip())
    ]


def _docx(path: Path) -> list[ExtractedUnit]:
    from docx import Document

    document = Document(str(path))
    paragraphs = "\n".join(p.text.strip() for p in document.paragraphs if p.text.strip())
    tables = []
    for table_number, table in enumerate(document.tables, start=1):
        rows = [" | ".join(cell.text.strip() for cell in row.cells) for row in table.rows]
        tables.append(ExtractedUnit("\n".join(rows), f"таблица {table_number}"))
    units = [ExtractedUnit(paragraphs, "текст")] if paragraphs else []
    return [*units, *tables]


def _docx_links(path: Path) -> tuple[DocumentLink, ...]:
    from docx import Document
    from docx.opc.constants import RELATIONSHIP_TYPE as RT
    from docx.oxml.ns import qn

    document = Document(str(path))
    links: list[DocumentLink] = []
    seen_targets: set[str] = set()
    for hyperlink in document.element.body.iter(qn("w:hyperlink")):
        relation_id = hyperlink.get(qn("r:id"))
        if not relation_id:
            continue
        relation = document.part.rels.get(relation_id)
        if relation is None or relation.reltype != RT.HYPERLINK or not relation.is_external:
            continue
        target = str(relation.target_ref).strip()
        if target in seen_targets or not _safe_external_target(target):
            continue
        label = " ".join(
            text.strip()
            for node in hyperlink.iter(qn("w:t"))
            if (text := node.text) and text.strip()
        )
        links.append(DocumentLink(label=(label or target)[:200], target=target))
        seen_targets.add(target)
        if len(links) >= 50:
            break
    return tuple(links)


def _safe_external_target(target: str) -> bool:
    if not target or len(target) > 4096 or any(ord(character) < 32 for character in target):
        return False
    try:
        parsed = urlsplit(target)
    except ValueError:
        return False
    scheme = parsed.scheme.casefold()
    if scheme in {"http", "https"}:
        return bool(parsed.netloc)
    if scheme == "mailto":
        return bool(parsed.path and "@" in parsed.path)
    return False


def _link_unit(links: tuple[DocumentLink, ...]) -> ExtractedUnit:
    entries = "\n".join(f"- {link.label}: {link.target}" for link in links)
    return ExtractedUnit(
        f"Внешние ссылки из документа (точные адреса):\n{entries}",
        "внешние ссылки",
    )


def _xlsx(path: Path) -> list[ExtractedUnit]:
    from openpyxl import load_workbook

    workbook = load_workbook(path, read_only=True, data_only=True)
    units: list[ExtractedUnit] = []
    for sheet in workbook.worksheets:
        rows = [
            " | ".join("" if value is None else str(value) for value in row)
            for row in sheet.iter_rows(values_only=True)
        ]
        text = "\n".join(row for row in rows if row.strip(" |"))
        if text:
            units.append(ExtractedUnit(text, f"лист «{sheet.title}»"))
    workbook.close()
    return units


def _pptx(path: Path) -> list[ExtractedUnit]:
    from pptx import Presentation

    presentation = Presentation(str(path))
    units: list[ExtractedUnit] = []
    for number, slide in enumerate(presentation.slides, start=1):
        text = "\n".join(
            shape.text.strip()
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        )
        if text:
            units.append(ExtractedUnit(text, f"слайд {number}"))
    return units


def _html(path: Path) -> list[ExtractedUnit]:
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(_decode(path.read_bytes()), "html.parser")
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()
    return [ExtractedUnit(soup.get_text("\n", strip=True), "страница")]


EXTRACTORS: dict[str, Callable[[Path], list[ExtractedUnit]]] = {
    ".txt": _text,
    ".md": _text,
    ".csv": _csv,
    ".pdf": _pdf,
    ".docx": _docx,
    ".xlsx": _xlsx,
    ".pptx": _pptx,
    ".html": _html,
    ".htm": _html,
}

SUPPORTED_EXTENSIONS = frozenset(EXTRACTORS)


def extract_document(path: Path) -> ExtractedDocument:
    suffix = path.suffix.lower()
    extractor = EXTRACTORS.get(suffix)
    if extractor is None:
        raise UnsupportedDocumentError(f"Формат {suffix or '<без расширения>'} не поддерживается")
    links = _docx_links(path) if suffix == ".docx" else ()
    units = tuple(unit for unit in extractor(path) if unit.text.strip())
    if links:
        units = (*units, _link_unit(links))
    if not units:
        raise EmptyDocumentError(f"В документе {path.name} нет извлекаемого текста")
    return ExtractedDocument(title=path.name, units=units, links=links)
