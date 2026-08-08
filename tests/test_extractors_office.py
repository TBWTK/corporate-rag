from pathlib import Path

import pytest
from docx import Document
from docx.opc.constants import RELATIONSHIP_TYPE as RT
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from openpyxl import Workbook
from pptx import Presentation
from reportlab.pdfgen import canvas

from rag_app.ingestion.extractors import DocumentLink, EmptyDocumentError, extract_document


def _add_hyperlink(paragraph: object, label: str, target: str) -> None:
    part = paragraph.part  # type: ignore[attr-defined]
    relation_id = part.relate_to(target, RT.HYPERLINK, is_external=True)
    hyperlink = OxmlElement("w:hyperlink")
    hyperlink.set(qn("r:id"), relation_id)
    run = OxmlElement("w:r")
    text = OxmlElement("w:t")
    text.text = label
    run.append(text)
    hyperlink.append(run)
    paragraph._p.append(hyperlink)  # type: ignore[attr-defined]


def test_extract_docx_with_paragraph_and_table(tmp_path: Path) -> None:
    path = tmp_path / "policy.docx"
    document = Document()
    document.add_paragraph("Правило удалённой работы")
    table = document.add_table(rows=1, cols=2)
    table.cell(0, 0).text = "Лимит"
    table.cell(0, 1).text = "2 дня"
    document.save(path)

    result = extract_document(path)

    assert [unit.location for unit in result.units] == ["текст", "таблица 1"]
    assert "Лимит | 2 дня" in result.units[1].text


def test_extract_docx_preserves_safe_external_hyperlinks(tmp_path: Path) -> None:
    path = tmp_path / "vpn.docx"
    document = Document()
    paragraph = document.add_paragraph("Скачайте архив: ")
    _add_hyperlink(paragraph, "VPN Client.zip", "https://downloads.example/vpn.zip")
    _add_hyperlink(paragraph, "Дубликат", "https://downloads.example/vpn.zip")
    _add_hyperlink(paragraph, "Поддержка", "mailto:helpdesk@example.com")
    _add_hyperlink(paragraph, "Опасная ссылка", "javascript:alert(1)")
    document.save(path)

    result = extract_document(path)

    assert result.links == (
        DocumentLink(label="VPN Client.zip", target="https://downloads.example/vpn.zip"),
        DocumentLink(label="Поддержка", target="mailto:helpdesk@example.com"),
    )
    link_unit = next(unit for unit in result.units if unit.location == "внешние ссылки")
    assert "VPN Client.zip: https://downloads.example/vpn.zip" in link_unit.text
    assert "Поддержка: mailto:helpdesk@example.com" in link_unit.text
    assert "javascript:" not in link_unit.text


def test_extract_xlsx_preserves_sheet_location(tmp_path: Path) -> None:
    path = tmp_path / "limits.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Лимиты"
    sheet.append(["Такси", 2500])
    workbook.save(path)

    result = extract_document(path)

    assert result.units[0].location == "лист «Лимиты»"
    assert "Такси | 2500" in result.units[0].text


def test_extract_pptx_by_slide(tmp_path: Path) -> None:
    path = tmp_path / "briefing.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Инструктаж"
    slide.placeholders[1].text = "Сообщите об инциденте за 15 минут"
    presentation.save(path)

    result = extract_document(path)

    assert result.units[0].location == "слайд 1"
    assert "15 минут" in result.units[0].text


def test_extract_pdf_by_page(tmp_path: Path) -> None:
    path = tmp_path / "notice.pdf"
    pdf = canvas.Canvas(str(path))
    pdf.drawString(72, 760, "Hotel limit is 10000 rubles")
    pdf.showPage()
    pdf.save()

    result = extract_document(path)

    assert result.units[0].location == "стр. 1"
    assert "10000" in result.units[0].text


def test_extract_html_removes_scripts(tmp_path: Path) -> None:
    path = tmp_path / "page.html"
    path.write_text(
        "<h1>Политика</h1><script>steal()</script><p>Используйте VPN</p>", encoding="utf-8"
    )

    result = extract_document(path)

    assert "Используйте VPN" in result.units[0].text
    assert "steal" not in result.units[0].text


def test_extract_cp1251_and_reject_empty(tmp_path: Path) -> None:
    encoded = tmp_path / "legacy.txt"
    encoded.write_bytes("Старый документ".encode("cp1251"))
    assert "Старый" in extract_document(encoded).units[0].text

    empty = tmp_path / "empty.md"
    empty.write_text("  \n", encoding="utf-8")
    with pytest.raises(EmptyDocumentError):
        extract_document(empty)
